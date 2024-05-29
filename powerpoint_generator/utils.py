import base64
import difflib
import imghdr
import io
import json
import logging
import os
import google.generativeai as genai
import requests
from PIL import Image, UnidentifiedImageError
from django.conf import settings
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


def get_font_colors(template_choice):
    if template_choice == 'dark_modern':
        return RGBColor(255, 255, 255), RGBColor(255, 255, 255)
    elif template_choice == 'bright_modern':
        return RGBColor(0, 0, 0), RGBColor(0, 0, 0)
    else:
        return RGBColor(0, 0, 0), RGBColor(0, 0, 0)


def apply_formatting(paragraphs, font_name, font_color, font_size):
    for paragraph in paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = font_size
            run.font.color.rgb = font_color


def extract_form_data(post_data):
    required_fields = ['topic_name', 'no_of_slides', 'user_role', 'tone', 'target_audience', 'prompt',
                       'template_choice']
    for field in required_fields:
        value = post_data.get(field, '').strip()
        if not value:
            return None
        elif field == 'no_of_slides' and not value.isdigit():
            return None
    return {
        'topic_name': post_data['topic_name'],
        'no_of_slides': int(post_data['no_of_slides']),
        'user_role': post_data['user_role'],
        'tone': post_data['tone'],
        'target_audience': post_data['target_audience'],
        'prompt': post_data['prompt'],
        'template_choice': post_data['template_choice'],
    }


def generate_prompt(form_data):
    payload = (
        f"Create a presentation on the topic '{form_data['topic_name']}'. You are acting as a '{form_data['user_role']}'"
        f"targeting an audience interested in '{form_data['target_audience']}'. The presentation should convey the "
        f"following information: '{form_data['prompt']}' and maintain a '{form_data['tone']}' tone throughout. The "
        f"presentation should be divided into {form_data['no_of_slides']} slides. "
        f"The response should be in JSON format as mentioned"
    )

    slides_content_json = json.dumps({
        "Slide 1": {
            "title": " ",
            "content": " ",
            "prompt_for_image_for_this_slide": " "
        },
        "Slide 2": {
            "title": " ",
            "content": " ",
            "prompt_for_image_for_this_slide": " "
        },
    }, indent=4)

    payload += slides_content_json

    return generate_content(payload)


def generate_content(payload):
    try:
        genai.configure(api_key=settings.API_KEY_GEMINI)
        model = genai.GenerativeModel('gemini-pro')
        response = model.generate_content(payload)
        return response.text
    except Exception as e:
        logger.error(f"Error generating content: {e}")
        return 'Failed to generate content'


def make_image_request(url: str, payload: dict) -> bytes:
    try:
        response = requests.post(url, headers=settings.HEADERS, json=payload)
        response.raise_for_status()
        return response.content
    except requests.RequestException as e:
        logger.error(f"Error fetching image from {url}: {e}")
        raise


def search_images(payload: dict) -> bytes:
    endpoints = [
        settings.STABLE_DIFFUSION_URL,
        settings.OPEN_JOURNEY_URL,
        settings.WAIFU_DIFFUSION_URL
    ]

    for endpoint in endpoints:
        try:
            return make_image_request(endpoint, payload)
        except requests.RequestException as e:
            print(f"Error with endpoint {endpoint}: {e}")
            continue

    raise Exception("All endpoints failed to process the request.")


def process_image_bytes(image_bytes: bytes) -> str:
    try:
        image_io = io.BytesIO(image_bytes)
        image = Image.open(image_io)
        if image.format not in ["JPEG", "PNG", "GIF"]:
            raise ValueError("Unsupported image format")
        image_url = f"data:image/{image.format.lower()};base64," + base64.b64encode(image_bytes).decode('utf-8')
        return image_url
    except (UnidentifiedImageError, ValueError) as e:
        logger.error(f"Error processing image bytes: {e}")
        raise


def find_most_similar_layout(prs, target_name):
    layout_names = [layout.name for layout in prs.slide_layouts]
    closest_matches = difflib.get_close_matches(target_name, layout_names)
    if closest_matches:
        closest_match = closest_matches[0]
        for layout in prs.slide_layouts:
            if layout.name == closest_match:
                return layout
    return None


def find_content_placeholder(slide):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.placeholder_format.idx == 1:
            return shape
    return None


def delete_all_slides(prs):
    # Delete all slides from the presentation
    slides_to_remove = prs.slides._sldIdLst[:]
    slide_ids_to_remove = [slide.rId for slide in prs.slides._sldIdLst]
    for slide_id in slide_ids_to_remove:
        prs.part.drop_rel(slide_id)
    for slide_id in slides_to_remove:
        prs.slides._sldIdLst.remove(slide_id)


def create_ppt(template_choice, slides_content, presentation_title):
    template_choice = template_choice or "dark_modern"
    template_path = os.path.join('static', 'stock', 'templates', f"{template_choice}.pptx")
    if not os.path.exists(template_path):
        logger.error(f"Template file {template_path} not found.")
        return None

    try:
        prs = Presentation(template_path)
        title_font_color, content_font_color = get_font_colors(template_choice)

        delete_all_slides(prs)

        # Find the appropriate layouts
        title_slide_layout = find_most_similar_layout(prs, "Title Slide")
        content_slide_layout = find_most_similar_layout(prs, "Picture with Caption")
        if not title_slide_layout or not content_slide_layout:
            logger.warning("Necessary layouts not found in the template.")
            return None

        # Title slide
        title_slide = prs.slides.add_slide(title_slide_layout)
        title = title_slide.shapes.title
        title.text = presentation_title
        apply_formatting(title.text_frame.paragraphs, font_name='Times New Roman', font_color=title_font_color,
                         font_size=Pt(45))

        # Remove the subtitle placeholder if it exists
        if len(title_slide.placeholders) > 1:
            subtitle_placeholder = title_slide.placeholders[1]
            sp = subtitle_placeholder._sp
            sp.getparent().remove(sp)

        # Validate and parse the JSON content
        try:
            slides_content_dict = json.loads(slides_content)
        except json.JSONDecodeError as e:
            logger.error(f"JSON decoding error: {e}")
            return None

        # Iterate over the slides content dictionary
        for slide_number, slide_content in slides_content_dict.items():
            slide = prs.slides.add_slide(content_slide_layout)

            title_placeholder = slide.placeholders[0]
            content_placeholder = slide.placeholders[2]

            # Set the title and content in the slide placeholders
            title_placeholder.text = slide_content['title']
            content_placeholder.text = slide_content['content']

            apply_formatting(title_placeholder.text_frame.paragraphs, font_name='Times New Roman',
                             font_color=title_font_color, font_size=Pt(24))
            apply_formatting(content_placeholder.text_frame.paragraphs, font_name='Times New Roman',
                             font_color=content_font_color, font_size=Pt(16))

            # Generate and insert image for the slide
            if 'prompt_for_image_for_this_slide' in slide_content:
                image_prompt = slide_content['prompt_for_image_for_this_slide']
                if image_prompt:
                    logger.info(f"Searching images for: {image_prompt}")
                    image_bytes = search_images({'inputs': image_prompt})
                    if image_bytes:
                        try:
                            image_extension = 'png'
                            if image_extension in ['jpeg', 'png', 'gif']:
                                image = Image.open(io.BytesIO(image_bytes))
                                output_image_path = os.path.join('static', 'stock', 'images',
                                                                 f"image_{slide_number}.{image_extension}")
                                image.save(output_image_path)

                                image_placeholder = slide.placeholders[1] if len(slide.placeholders) > 2 else None
                                if image_placeholder:
                                    image_placeholder.insert_picture(output_image_path)
                                else:
                                    slide.shapes.add_picture(output_image_path, Inches(1), Inches(1))
                                os.remove(output_image_path)  # Remove the temporary image file
                            else:
                                logger.error("Unsupported image format")
                        except Exception as e:
                            logger.error(f"Error inserting image: {e}")

        output_dir = os.path.join('static', 'stock', 'presentations')
        os.makedirs(output_dir, exist_ok=True)

        sanitized_title = "".join(c for c in presentation_title if c.isalnum() or c in [' ', '_', '-'])
        output_path = os.path.join(output_dir, f"{sanitized_title}.pptx")
        prs.save(output_path)
        return output_path
    except Exception as e:
        logger.error(f"Error creating presentation: {e}")
        return None
